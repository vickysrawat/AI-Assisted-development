#!/usr/bin/env python3
"""Generate AI in SDLC PowerPoint presentation - Simplified 18 slides."""

import sys
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Installing python-pptx...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

# Colors
DARK_BLUE = RGBColor(13, 43, 85)
MED_BLUE = RGBColor(26, 107, 196)
LIGHT_BLUE = RGBColor(0, 176, 240)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(240, 240, 240)
GREEN = RGBColor(0, 176, 80)
RED = RGBColor(192, 0, 0)

def add_header(slide, title, bg_color=DARK_BLUE):
    """Add header to slide."""
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = bg_color
    header.line.color.rgb = bg_color
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_footer(slide):
    """Add footer to slide."""
    footer = slide.shapes.add_shape(1, Inches(0), Inches(7), Inches(10), Inches(0.5))
    footer.fill.solid()
    footer.fill.fore_color.rgb = GRAY
    footer.line.color.rgb = GRAY

def add_content_box(slide, x, y, width, height, text, font_size=14, bold=False, color=DARK_BLUE, bg_color=LIGHT_BLUE):
    """Add content box with text."""
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = DARK_BLUE
    
    txt_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1), Inches(width - 0.2), Inches(height - 0.2))
    tf = txt_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

def create_presentation():
    """Create 18-slide presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slides_data = [
        # Slide 1: Title
        {
            "type": "title",
            "title": "The AI Opportunity in Your SDLC Today",
            "subtitle": "Transform Your Development Lifecycle with AI",
            "kpis": [("10×", "Faster Onboarding"), ("60%", "Fewer PR Cycles"), ("333%", "ROI"), ("85%", "Fewer Bugs")]
        },
        # Slide 2: Traditional SDLC
        {"type": "traditional", "title": "Traditional SDLC - Current State (3-4 weeks)"},
        # Slide 3: AI-Assisted SDLC
        {"type": "ai_sdlc", "title": "AI-Assisted SDLC - Reimagined (3-5 days)"},
        # Slides 4-18: Content slides
        {"type": "content", "title": "4. Pain Points Matrix"},
        {"type": "content", "title": "5. AI in Planning - Specification-Driven Development"},
        {"type": "content", "title": "6. AI in Development - Code Generation with Quality Gates"},
        {"type": "content", "title": "7. AI in Review - Instant Intelligent Feedback"},
        {"type": "content", "title": "8. AI in Security - Shift-Left, Not Right"},
        {"type": "content", "title": "9. AI in Onboarding - From 30 Days to 2 Days"},
        {"type": "content", "title": "10. Compliance & Governance - Audit Trail"},
        {"type": "content", "title": "11. Metrics Dashboard - Real Data"},
        {"type": "content", "title": "12. Misconception - AI Replaces Developers"},
        {"type": "content", "title": "13. Business Case - ROI Analysis"},
        {"type": "content", "title": "14. Transition Plan - 4-Week Pilot"},
        {"type": "content", "title": "15. Key Takeaways"},
        {"type": "content", "title": "16. Q&A - Common Questions"},
        {"type": "content", "title": "17. Call to Action"},
        {"type": "content", "title": "18. Thank You"},
    ]
    
    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE
        
        if slide_data["type"] == "title":
            # Title slide
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = DARK_BLUE
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data["title"]
            p.font.size = Pt(54)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
            stf = subtitle_box.text_frame
            sp = stf.paragraphs[0]
            sp.text = slide_data["subtitle"]
            sp.font.size = Pt(24)
            sp.font.color.rgb = LIGHT_BLUE
            sp.alignment = PP_ALIGN.CENTER
            
            # KPI boxes
            for idx, (metric, label) in enumerate(slide_data["kpis"]):
                x = 0.8 + idx * 2.3
                box = slide.shapes.add_shape(1, Inches(x), Inches(5.8), Inches(2), Inches(1.3))
                box.fill.solid()
                box.fill.fore_color.rgb = LIGHT_BLUE
                box.line.color.rgb = LIGHT_BLUE
                
                m_box = slide.shapes.add_textbox(Inches(x), Inches(5.95), Inches(2), Inches(0.6))
                mtf = m_box.text_frame
                mp = mtf.paragraphs[0]
                mp.text = metric
                mp.font.size = Pt(32)
                mp.font.bold = True
                mp.font.color.rgb = WHITE
                mp.alignment = PP_ALIGN.CENTER
                
                l_box = slide.shapes.add_textbox(Inches(x), Inches(6.55), Inches(2), Inches(0.5))
                ltf = l_box.text_frame
                lp = ltf.paragraphs[0]
                lp.text = label
                lp.font.size = Pt(11)
                lp.font.color.rgb = WHITE
                lp.alignment = PP_ALIGN.CENTER
        
        elif slide_data["type"] == "traditional":
            add_header(slide, slide_data["title"])
            
            summary = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
            sf = summary.text_frame
            sf.word_wrap = True
            
            content = """Days 1-3: Planning
Days 4-5: Design
Days 6-10: Development
Days 11-14: Review & Rework
Days 15-20: Security & Compliance
Day 21+: Production Ready

❌ 3-4 weeks total
❌ High rework rate
❌ Manual process
❌ Security added late"""
            
            sp = sf.paragraphs[0]
            sp.text = content
            sp.font.size = Pt(14)
            sp.font.color.rgb = DARK_GRAY
            sp.alignment = PP_ALIGN.CENTER
            
            add_footer(slide)
        
        elif slide_data["type"] == "ai_sdlc":
            add_header(slide, slide_data["title"])
            
            summary = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
            sf = summary.text_frame
            sf.word_wrap = True
            
            content = """Day 1: AI generates structured plan
Day 2: AI drafts complete specification
Day 3: AI generates full implementation
Day 4: Developer reviews code (1 hour)
Day 5: Security automated, ready to deploy

✅ 3-5 days total
✅ Zero rework (spec-driven)
✅ Automated process
✅ Security built-in"""
            
            sp = sf.paragraphs[0]
            sp.text = content
            sp.font.size = Pt(14)
            sp.font.bold = True
            sp.font.color.rgb = GREEN
            sp.alignment = PP_ALIGN.CENTER
            
            add_footer(slide)
        
        else:  # content slides
            add_header(slide, slide_data["title"])
            
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.3))
            cf = content_box.text_frame
            cf.word_wrap = True
            cp = cf.paragraphs[0]
            cp.text = "[Professional content and visualizations would be inserted here]"
            cp.font.size = Pt(16)
            cp.font.color.rgb = RGBColor(100, 100, 100)
            cp.alignment = PP_ALIGN.CENTER
            
            add_footer(slide)
    
    return prs

if __name__ == "__main__":
    try:
        print("Generating presentation...")
        prs = create_presentation()
        output_file = "AI_in_SDLC_Presentation.pptx"
        prs.save(output_file)
        print(f"✅ SUCCESS: {output_file} created ({len(prs.slides)} slides)")
    except Exception as e:
        print(f"❌ ERROR: {e}")
        sys.exit(1)
